#!/usr/bin/env python3
"""Fill assets/template.pptx from a deck_spec.json.

Usage:
    python3 fill.py deck_spec.json out.pptx

deck_spec.json shape:
{
  "drop_slides": [9, 20],          // 1-based slide indices to remove, e.g. the
                                    // POV-plan slide for a post-POC review deck
  "remove_slots": {
    "6": ["NEED_5"],                // slide 6, remove the NEED_5 title+desc shapes
    "17": ["STEP_3", "STEP_4"]      // slide 17 (next steps), remove unused rows
  },
  "values": {
    "CLIENT": "Astra Security",
    "CHAMPION_GOAL": "...",
    ...
    "ACV_TOTAL": "[[TBD: commercials not yet agreed]]"
  }
}

Rules:
- Every token in the template must have a value in `values`, UNLESS its slot
  is listed in `remove_slots` (see below) or its slide is in `drop_slides`.
- A value may itself be a `[[TBD: ...]]` placeholder — that's allowed, it just
  gets reported at the end so a human fills it before the deck ships.
- `remove_slots` removes a token's enclosing shape (or table row, if the token
  is in a table cell) instead of filling it — used for slots the source decks
  show as optional (e.g. NEED_5 on slide 6, or a commercials line item).
  The slot key is the token's shared prefix (e.g. "NEED_5" covers both
  NEED_5_TITLE and NEED_5_DESC).
"""
import copy
import json
import re
import sys
from pptx import Presentation

TOKEN_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")


def iter_text_frames_with_removable(shapes):
    """Yield (text_frame, remove_fn) pairs. remove_fn() deletes the shape/row
    that text_frame lives in."""
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_text_frames_with_removable(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    def remove_row(row=row, tbl=shape.table):
                        tbl._tbl.remove(row._tr)
                    yield cell.text_frame, remove_row
            continue
        if shape.has_text_frame:
            def remove_shape(shape=shape):
                shape._element.getparent().remove(shape._element)
            yield shape.text_frame, remove_shape


def slot_of(token):
    """NEED_5_TITLE -> NEED_5, STEP_3_TASK -> STEP_3, CLIENT -> CLIENT."""
    m = re.match(r"^([A-Z]+_\d+)(_.*)?$", token)
    return m.group(1) if m else token


def fill(template_path, spec, out_path):
    prs = Presentation(template_path)

    drop = set(spec.get("drop_slides", []))
    remove_slots = {int(k): set(v) for k, v in spec.get("remove_slots", {}).items()}
    values = spec.get("values", {})

    if drop:
        ids = list(prs.slides._sldIdLst)
        to_remove = [ids[i - 1] for i in sorted(drop)]
        for el in to_remove:
            prs.slides._sldIdLst.remove(el)

    # remove_slots' slide-number keys are purely organizational for spec
    # authors — token slot names (e.g. "NEED_5") are already unique across
    # the whole template, so matching is done by slot name only, below.
    missing = []
    tbds = []

    for slide in prs.slides:
        for tf, remove_fn in list(iter_text_frames_with_removable(slide.shapes)):
            for p in tf.paragraphs:
                text = "".join(r.text for r in p.runs)
                tokens = TOKEN_RE.findall(text)
                if not tokens:
                    continue
                slot_removed = False
                for tok in tokens:
                    slot = slot_of(tok)
                    for removed_slots in remove_slots.values():
                        if slot in removed_slots:
                            slot_removed = True
                if slot_removed:
                    remove_fn()
                    continue
                new_text = text
                for tok in tokens:
                    if tok not in values:
                        missing.append(tok)
                        continue
                    val = str(values[tok])
                    if val.startswith("[[TBD"):
                        tbds.append((tok, val))
                    new_text = new_text.replace("{{%s}}" % tok, val)
                if new_text != text and p.runs:
                    p.runs[0].text = new_text
                    for r in p.runs[1:]:
                        r.text = ""

    if missing:
        print("ERROR: missing values for tokens:", sorted(set(missing)), file=sys.stderr)
        sys.exit(1)

    prs.save(out_path)
    print(f"wrote {out_path}")
    if tbds:
        print(f"\n{len(tbds)} TBD placeholder(s) left in the deck — fill before sending:")
        for tok, val in tbds:
            print(f"  {tok}: {val}")


def main():
    if len(sys.argv) < 3:
        print("usage: fill.py deck_spec.json out.pptx", file=sys.stderr)
        sys.exit(2)
    spec_path, out_path = sys.argv[1], sys.argv[2]
    with open(spec_path) as f:
        spec = json.load(f)
    import os
    template = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets", "template.pptx")
    fill(template, spec, out_path)


if __name__ == "__main__":
    main()
