#!/usr/bin/env python3
"""Self-check for fill.py. assert-based, no framework.

Run: python3 test_fill.py
"""
import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from fill import fill

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "..", "assets", "template.pptx")
SPEC = os.path.join(HERE, "..", "examples", "astra_deck_spec.json")
OUT = "/tmp/poc_deck_test_out.pptx"


def all_text(path):
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    def walk(shapes):
        for shp in shapes:
            if shp.shape_type == 6:
                walk(shp.shapes); continue
            if getattr(shp, "has_table", False) and shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            chunks.append("".join(r.text for r in p.runs))
                continue
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    chunks.append("".join(r.text for r in p.runs))
    for s in prs.slides:
        walk(s.shapes)
    return "\n".join(chunks), len(prs.slides)


def test_full_fill_no_tokens_left():
    with open(SPEC) as f:
        spec = json.load(f)
    fill(TEMPLATE, spec, OUT)
    text, n_slides = all_text(OUT)
    assert "{{" not in text, "unfilled token survived into output"
    assert "Suraj" in text, "champion name missing from filled deck"
    assert n_slides == 26, f"expected 26 slides (28 - 2 dropped), got {n_slides}"
    print("ok: full fill, zero leftover tokens, champion name present, slide count correct")


def test_remove_slot_deletes_shape():
    with open(SPEC) as f:
        spec = json.load(f)
    spec = dict(spec)
    spec["remove_slots"] = {"6": ["NEED_5"]}
    fill(TEMPLATE, spec, OUT)
    text, _ = all_text(OUT)
    assert "One bar, every rep" not in text, "NEED_5 slot should have been removed"
    print("ok: remove_slots deletes the targeted shape")


def test_missing_token_hard_errors():
    with open(SPEC) as f:
        spec = json.load(f)
    spec = dict(spec)
    spec["values"] = dict(spec["values"])
    del spec["values"]["CLIENT"]
    try:
        fill(TEMPLATE, spec, OUT)
        assert False, "expected SystemExit for missing token"
    except SystemExit as e:
        assert e.code != 0
    print("ok: missing token hard-errors")


if __name__ == "__main__":
    test_full_fill_no_tokens_left()
    test_remove_slot_deletes_shape()
    test_missing_token_hard_errors()
    print("\nall tests passed")
