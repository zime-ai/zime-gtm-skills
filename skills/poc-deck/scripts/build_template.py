#!/usr/bin/env python3
"""One-off script: build assets/template.pptx from the two sample decks.

Not part of the skill's runtime — kept for provenance and for rebuilding the
template if the source decks are revised. Requires the original two sample
decks on disk (not part of this repo); update ASTRA/TRUEFOUNDRY below to
wherever they live before rerunning.

Steps:
  1. Start from the Astra deck (16:9, superset of slide types).
  2. Port in the 2 slides Astra lacks but TrueFoundry has (KPI funnel,
     3-week POV plan), scaled from TrueFoundry's 10x5.62in canvas to
     Astra's 13.33x7.5in canvas, and move them to sensible positions.
  3. Replace every client-specific string with a {{TOKEN}}, leaving fixed
     boilerplate copy untouched.
"""
import copy
import os
import sys
from pptx import Presentation
from pptx.oxml.ns import qn

ASTRA = "/Users/atharvadiwan/Antigravity_workspace/Master_workspace/zime_repos/sample_poc_deck/Zime_Astra_D7_AUG 1.pptx"
TRUEFOUNDRY = "/Users/atharvadiwan/Antigravity_workspace/Master_workspace/zime_repos/sample_poc_deck/TrueFoundry_Post_POC.pptx"
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets", "template.pptx")


def scale_shape_xml(el, factor):
    """Scale every a:off / a:ext x/y/cx/cy found under el by factor in place."""
    for off in el.iter(qn("a:off")):
        off.set("x", str(round(int(off.get("x")) * factor)))
        off.set("y", str(round(int(off.get("y")) * factor)))
    for ext in el.iter(qn("a:ext")):
        ext.set("cx", str(round(int(ext.get("cx")) * factor)))
        ext.set("cy", str(round(int(ext.get("cy")) * factor)))


def copy_slide(dest_prs, src_slide, scale_factor=1.0):
    layout = dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src_slide.shapes:
        new_el = copy.deepcopy(shp._element)
        if scale_factor != 1.0:
            scale_shape_xml(new_el, scale_factor)
        new_slide.shapes._spTree.append(new_el)
    return new_slide


def move_slides(prs, moves):
    """moves: list of (from_end_offset, to_index) applied by object identity."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for obj, to_index in moves:
        ids.remove(obj)
        ids.insert(to_index, obj)
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in ids:
        sldIdLst.append(el)


# ---------------------------------------------------------------------------
# Global name substitutions — applied to any paragraph not in CONTENT_TOKEN_MAP.
# Order matters: longer/more specific patterns first.
# ---------------------------------------------------------------------------
GLOBAL_SUBS = [
    ("Astra Security", "{{CLIENT}}"),
    ("ASTRA SECURITY", "{{CLIENT_UPPER}}"),
    ("Astra’s", "{{CLIENT_SHORT}}’s"),
    ("ASTRA’S", "{{CLIENT_SHORT_UPPER}}’S"),
    ("Astra", "{{CLIENT_SHORT}}"),
    ("Suraj’s", "{{CHAMPION_FIRST}}’s"),
    ("Suraj", "{{CHAMPION_FIRST}}"),
    ("TrueFoundry", "{{CLIENT}}"),
]

# ---------------------------------------------------------------------------
# Content token map — exact original paragraph text -> token/replacement.
# Applied BEFORE global subs, on an exact match.
# ---------------------------------------------------------------------------
CONTENT_TOKEN_MAP = {
    # Slide 2 — champion goal
    "To have a Inbound system that will get to 30% conversion on 350 MQLs per month": "{{CHAMPION_GOAL}}",

    # Slide 3 — blockers (fixed at 3 slots; both sample decks use exactly 3)
    "Inconsistent Qualification": "{{BLOCKER_1_TITLE}}",
    "No shared standard, judgment lives with a couple of senior reps instead of a documented system": "{{BLOCKER_1_DESC}}",
    "Lack of Visibility on Deal Health": "{{BLOCKER_2_TITLE}}",
    "No reliable read on pipeline without re-listening to hours of calls": "{{BLOCKER_2_DESC}}",
    "No Action Layer": "{{BLOCKER_3_TITLE}}",
    "Insights stay buried in the recording; nothing routes back to reps as actionable next steps - based on past wins": "{{BLOCKER_3_DESC}}",

    # Slide 4 — root cause
    "Your best practices for Qualification ": "Your best practices for {{INITIATIVE_NAME}} ",
    "FAINT Qualification": "{{RUBRIC_NAME}}",
    "Top reps head ": "{{ROOT_CAUSE_1}} ",
    "Hours of MeetGeek recordings ": "{{ROOT_CAUSE_2}} ",

    # Slide 5 — impact / risk tiles (fixed at 4 slots)
    "Risk of Inconsistent execution in a high-velocity cycle": "{{RISK_HEADLINE}}",
    "REVENUE": "{{RISK_1_LABEL}}",
    "30-35% of potential wins lost to inconsistent qualification": "{{RISK_1_STAT}}",
    "High-velocity team lose deals at discovery and qualification without consistent best practices. ": "{{RISK_1_DESC}}",
    "HIDDEN RISKS & STALL SIGNALS": "{{RISK_2_LABEL}}",
    "20–30% of pipeline quietly stalls, unqualified": "{{RISK_2_STAT}}",
    "Deals that were never really moving keep consuming rep time until it’s too late.": "{{RISK_2_DESC}}",
    "LEADERSHIP BECOMES BOTTLENECK": "{{RISK_3_LABEL}}",
    "4 hours of calls, per rep, per day, on Suraj alone to review": "{{RISK_3_STAT}}",
    "Without trusted signals, that review falls on leaders instead of the team running itself.": "{{RISK_3_DESC}}",
    "NEXT HIRE RISK": "{{RISK_4_LABEL}}",
    "Every new rep starts from a blank slate": "{{RISK_4_STAT}}",
    "Qualification instinct lives with a handful of seasoned reps. Without a standard playbook, everyone has to find their own way, and every new hire repeats that same slow climb.": "{{RISK_4_DESC}}",

    # Slide 6 — champion quote + immediate needs (optional slide; drop if unused)
    "“I want to run inbound like a machine”": "{{CHAMPION_QUOTE}}",
    "Suraj  ·  Director of Sales, Astra Security": "{{CHAMPION_FIRST}}  ·  {{CHAMPION_TITLE}}, {{CLIENT}}",
    "See every deal's true health": "{{NEED_1_DESC}}",
    "Numbers you can commit to": "{{NEED_2_DESC}}",
    "Minutes, not hours": "{{NEED_3_DESC}}",
    "Systemize before adding reps": "{{NEED_4_DESC}}",
    "One bar, every rep": "{{NEED_5_DESC}}",
    "Pipeline Visibility": "{{NEED_1_TITLE}}",
    "Forecast Accuracy": "{{NEED_2_TITLE}}",
    "Less Manual Review": "{{NEED_3_TITLE}}",
    "Structure Before Scaling": "{{NEED_4_TITLE}}",
    "Consistent Qualification": "{{NEED_5_TITLE}}",

    # Slide 8 — introducing (step 2 title is motion-dependent)
    "Step 2 - Apply it to every deal": "{{STEP2_TITLE}}",
    "Learns what makes you winat Discovery & Qualification": "{{STEP1_SUBTITLE}}",

    # Slide 10 — operationalization metrics (fixed at 4 slots)
    "10X": "{{METRIC_PRODUCTIVITY}}",
    "+20%": "{{METRIC_TIME_SAVED}}",
    "10%": "{{METRIC_HEADCOUNT_SAVED}}",
    "+18%": "{{METRIC_WIN_RATE}}",

    # Slide 12 — cost of inaction
    "$490k  in qualified pipeline goes unrealized every month": "{{ROI_HEADLINE_STAT}}",
    "Every month without a system, the inconsistency scales as the team does": "{{ROI_HEADLINE_DESC}}",

    # Slide 16 — commercials (fixed at 2 required lines + 1 optional 3rd line)
    "$": "{{ACV_TOTAL}}",
    "06 SME licenses": "{{LINE_1_LABEL}}",
    "$7,182 / yr": "{{LINE_1_PRICE}}",
    "$99 / user / month · Pipeline review, deal health dashboard, prep notes, and call insights. Includes 15 free credits for advanced AI features; after that $1/credit.": "{{LINE_1_DESC}}",
    "Annual platform fee": "{{LINE_2_LABEL}}",
    "$5,000 / yr": "{{LINE_2_PRICE}}",
    "Infrastructure, integrations, and forward-deployed engineer support. 1st qualification playbook is free, then $12K/playbook.": "{{LINE_2_DESC}}",
    "$12,182 / yr": "{{ACV_TOTAL_LINE}}",

    # Slide 18 — next steps (fixed at 4 slots)
    "Confirm onboarding date ": "{{STEP_1_TASK}}",
    "Loop in HubSpot admin to scope auto-field sync": "{{STEP_2_TASK}}",
    "Sign off SOW & MSA": "{{STEP_3_TASK}}",
    "Pilot User details": "{{STEP_4_TASK}}",
    "TBC": "{{STEP_ETA_TBC}}",
    "click here": "{{STEP_3_LINK}}",
    "Click here": "{{STEP_4_LINK}}",

    # Slide 19 — current stack comparison table
    "MeetGeek": "{{STACK_TOOL_1}}",
    "Hubspot": "{{STACK_TOOL_2}}",

    # ---- Ported from TrueFoundry: KPI funnel slide ----
    "Your Data  —  Pylon, Fireflies, Salesforce": "{{CLIENT_TOOLS}}",
    "Your Insights  —  ClaudeHard to drive adoption across org": "{{KPI_INSIGHTS_LINE}}",
    "Build Your Best PracticesKPI: 100+ best practices with 80%+ win correlation (e.g. BU intros, QBRs, etc.)": "{{KPI_BUILD_LINE}}",
    "Account IntelligenceKPI: 30% leader time saved with visibility of your best practices": "{{KPI_ACCOUNT_INTEL_LINE}}",
    "Your InitiativesKPI:  Expansion $$ (Other e.gs: discovery, churn, close rates)": "{{KPI_INITIATIVES_LINE}}",
    "Rep Coaching, Ramp UpKPI:  80% A-reps with adoption of your best practices": "{{KPI_REP_COACHING_LINE}}",

    # ---- Ported from TrueFoundry: 3-week POV plan slide ----
    "Kickoff & Data Pull": "{{POV_W1_TASK1_TITLE}}",
    "Recordings & deal names of 10 won deals": "{{POV_W1_TASK1_DESC}}",
    "Build Playbook v1": "{{POV_W1_TASK2_TITLE}}",
    "Correlate winning behaviors into a draft": "{{POV_W1_TASK2_DESC}}",
    "Deploy & Flag Risk": "{{POV_W2_TASK1_TITLE}}",
    "Apply playbook to 15 stalled deals": "{{POV_W2_TASK1_DESC}}",
    "Validate on Lost Deals": "{{POV_W2_TASK2_TITLE}}",
    "Apply to 10 lost deals, confirm loss patterns": "{{POV_W2_TASK2_DESC}}",
    "Pipeline Review Report": "{{POV_W3_TASK1_TITLE}}",
    "Compile findings across all 35 deals": "{{POV_W3_TASK1_DESC}}",
    "Demo & Go/No-Go Call": "{{POV_W3_TASK2_TITLE}}",
    "Go or no-go decision": "{{POV_W3_TASK2_DESC}}",
    "Accurate expansion signals & customer objections per deal": "{{POV_OUTCOME_1}}",
    "Helping CS and managers fix real signals and lift expansions.": "{{POV_OUTCOME_2}}",
}


def iter_text_frames(shapes):
    """Recurse into groups and tables, yield every text_frame."""
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_text_frames(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
            continue
        if shape.has_text_frame:
            yield shape.text_frame


def tokenize_paragraph(paragraph):
    runs = paragraph.runs
    if not runs:
        return
    original = "".join(r.text for r in runs)
    if not original:
        return
    if original in CONTENT_TOKEN_MAP:
        new_text = CONTENT_TOKEN_MAP[original]
    else:
        new_text = original
        for old, new in GLOBAL_SUBS:
            new_text = new_text.replace(old, new)
    if new_text != original:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""


def tokenize_presentation(prs):
    for slide in prs.slides:
        for tf in iter_text_frames(slide.shapes):
            for p in tf.paragraphs:
                tokenize_paragraph(p)


def main():
    astra = Presentation(ASTRA)
    tf = Presentation(TRUEFOUNDRY)

    scale = astra.slide_width / tf.slide_width  # uniform 4/3 scale, both dims match
    assert abs(scale - astra.slide_height / tf.slide_height) < 1e-6, "non-uniform scale, check canvas sizes"

    kpi_src = tf.slides[6]   # "You have data & insights. You need outcomes."
    pov_src = tf.slides[18]  # "POV: Let Us Show You Results"

    copy_slide(astra, kpi_src, scale_factor=scale)
    copy_slide(astra, pov_src, scale_factor=scale)

    # Both new slides were appended at the end; move them into place by
    # matching the last two sldId XML elements.
    ids = list(astra.slides._sldIdLst)
    kpi_id_el = ids[-2]
    pov_id_el = ids[-1]
    move_slides(astra, [(kpi_id_el, 8), (pov_id_el, 19)])

    tokenize_presentation(astra)

    astra.save(OUT)
    print(f"wrote {OUT} — {len(astra.slides)} slides")


if __name__ == "__main__":
    main()
