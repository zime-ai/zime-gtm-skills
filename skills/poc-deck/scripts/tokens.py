#!/usr/bin/env python3
"""List every {{TOKEN}} in a template.pptx, grouped by slide index (0-based).

Usage: python3 tokens.py [path/to/template.pptx]

This is the single source of truth for the deck_spec.json schema — fill.py
and any spec author should derive the required keys from this, not from a
hand-maintained list.
"""
import os
import re
import sys
from pptx import Presentation

TOKEN_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")
DEFAULT_TEMPLATE = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "assets", "template.pptx")


def iter_text_frames(shapes):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_text_frames(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
            continue
        if shape.has_text_frame:
            yield shape.text_frame


def find_tokens(path):
    prs = Presentation(path)
    by_slide = {}
    all_tokens = set()
    for i, slide in enumerate(prs.slides, 1):
        found = set()
        for tf in iter_text_frames(slide.shapes):
            for p in tf.paragraphs:
                text = "".join(r.text for r in p.runs)
                found.update(TOKEN_RE.findall(text))
        if found:
            by_slide[i] = sorted(found)
            all_tokens.update(found)
    return by_slide, sorted(all_tokens)


def main():
    path = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_TEMPLATE
    by_slide, all_tokens = find_tokens(path)
    for slide_no, toks in by_slide.items():
        print(f"slide {slide_no}: " + ", ".join(toks))
    print(f"\n{len(all_tokens)} unique tokens total:")
    for t in all_tokens:
        print(f"  {t}")


if __name__ == "__main__":
    main()
