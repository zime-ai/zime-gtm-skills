#!/usr/bin/env python3
"""Dump a .pptx into a diffable JSON structure: per slide, per shape, its
text/runs plus position and color -- so two decks (ground truth vs a
skill's generated output) can be compared without rendering either one.

No pixel rendering here (no libreoffice on this machine) -- this is
structural conformity only: text content, shape position/size, fill color,
font color/size/bold. A layout bug that doesn't touch any of these
properties (e.g. two boxes overlapping despite both having "correct"
coordinates) won't surface. See references/known-pitfalls.md.

Usage: extract_pptx.py <deck.pptx> [> out.json]
"""
import json
import sys

from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400


def emu_to_in(v):
    return round(v / EMU_PER_INCH, 3) if v is not None else None


def rgb_of(color_format):
    """Return 'RRGGBB' if this shape/run has an explicit solid RGB color,
    None otherwise (theme color, no fill, inherited -- anything not a
    literal RGB value, which is the only kind worth diffing directly)."""
    try:
        if color_format.type is not None and hasattr(color_format, "rgb"):
            return str(color_format.rgb)
    except (AttributeError, TypeError, KeyError):
        pass
    return None


def extract_runs(text_frame):
    runs = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            runs.append({
                "text": run.text,
                "font_size_pt": run.font.size.pt if run.font.size else None,
                "bold": run.font.bold,
                "color": rgb_of(run.font.color),
            })
    return runs


def extract_shape(shape):
    entry = {
        "name": shape.name,
        "position_in": {
            "left": emu_to_in(shape.left), "top": emu_to_in(shape.top),
            "width": emu_to_in(shape.width), "height": emu_to_in(shape.height),
        },
    }
    fill_color = None
    try:
        if shape.fill.type is not None:
            fill_color = rgb_of(shape.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        pass
    entry["fill_color"] = fill_color

    if getattr(shape, "has_table", False) and shape.has_table:
        entry["table"] = [
            [{"text": cell.text_frame.text, "runs": extract_runs(cell.text_frame)}
             for cell in row.cells]
            for row in shape.table.rows
        ]
    elif shape.shape_type == 6:  # GROUP
        entry["group"] = [extract_shape(s) for s in shape.shapes]
    elif shape.has_text_frame:
        entry["text"] = shape.text_frame.text
        entry["runs"] = extract_runs(shape.text_frame)

    return entry


def main():
    if len(sys.argv) != 2:
        sys.exit(f"Usage: {sys.argv[0]} <deck.pptx>")
    prs = Presentation(sys.argv[1])
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slides.append({"slide": i, "shapes": [extract_shape(s) for s in slide.shapes]})
    json.dump(slides, sys.stdout, indent=2)
    print()


if __name__ == "__main__":
    main()
